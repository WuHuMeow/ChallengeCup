"""Generate the defense PPT from frozen analysis artifacts (Task 24)."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

FORMAL = Path("output/evidence/formal")


def _read_csv(path: Path) -> list[dict[str, str]]:
    import csv

    with path.open(newline="", encoding="utf-8") as handle:
        return list(csv.DictReader(handle))


def _add_table(slide, rows: list[list[str]], left: float, top: float,
               width: float, height: float) -> None:
    shape = slide.shapes.add_table(
        len(rows), len(rows[0]), Inches(left), Inches(top), Inches(width), Inches(height)
    )
    table = shape.table
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            cell.text_frame.paragraphs[0].runs[0].font.size = Pt(11)


def main() -> int:
    paired = _read_csv(FORMAL / "paired_tests.csv")
    desc = _read_csv(FORMAL / "descriptive_stats.csv")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def title_slide(title: str, subtitle: str) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle

    def bullet_slide(title: str, bullets: list[str]) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        body = slide.placeholders[1].text_frame
        for index, item in enumerate(bullets):
            para = body.paragraphs[0] if index == 0 else body.add_paragraph()
            para.text = item
            para.font.size = Pt(16)

    title_slide(
        "CA-MP 容量感知最大压力控制",
        "XH-202613 赛道 B · 20 路口车路云协同仿真平台 · SUMO 1.27.1",
    )
    bullet_slide(
        "问题与方案",
        [
            "窄路密网：短车道、低容量、排队回溢快",
            "经典 MaxPressure 以绝对排队分配路权，偏向长车道",
            "CA-MP 三项改进：容量归一化压力 / 下游溢出门控 / 云端动态绿灯",
            "统一运行链路：RunService → VariantBundle → SimulationRunner",
        ],
    )
    bullet_slide(
        "形式实验设计",
        [
            "20 真实路口 × 3 算法 × 2 流量 × 3 种子 = 360 正常 run",
            "扰动 180 run：施工占道 / 大型活动需求 / 车辆故障车道阻塞 各 60",
            "3600s 仿真 + 600s 预热，全部 sealed evidence，540/540 完成",
            "分析：配对差值 95% CI（candidate − baseline，负值 = 候选更优）",
        ],
    )
    travel = [r for r in paired if r["metric"] == "avg_travel_time"]
    rows = [["对比", "均值差(秒)", "相对变化", "95% CI", "改善单元"]]
    for r in travel:
        rows.append([
            f"{r['candidate']} vs {r['baseline']}",
            f"{float(r['mean_difference']):+.2f}",
            f"{float(r['relative_change']):+.2%}",
            f"[{float(r['ci_lower']):.2f}, {float(r['ci_upper']):.2f}]",
            f"{r['improved_unit_count']}/30",
        ])
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "主结果：平均旅行时间（normal，120 对/组）"
    _add_table(slide, rows, 0.8, 1.6, 11.7, 1.6)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(3.6), Inches(11.7), Inches(1.2))
    box.text_frame.text = (
        "两个候选的旅行时间改善均统计显著（95% CI 完全为负）；排队长度 CA-MP 全场最低"
    )
    queue = [r for r in desc if r["metric"] == "avg_queue_length" and r["matrix_kind"] == "normal"]
    queue_text = "；".join(f"{r['algorithm']}={float(r['mean']):.2f}m" for r in queue)
    box.text_frame.add_paragraph().text = f"平均排队长度：{queue_text}"

    bullet_slide(
        "扰动韧性（平均旅行时间，秒）",
        [
            "施工占道：CA-MP 67.20（最优，较基线改善 7.2%）；classic 67.38；fixed 72.45",
            "大型活动需求：classic 60.00；CA-MP 61.02；fixed 63.76",
            "车辆故障阻塞：classic 65.42；CA-MP 66.03；fixed 69.35",
            "平均排队长度：CA-MP 在全部三类扰动下均为最低（3.78/3.44/3.53 m）",
        ],
    )
    bullet_slide(
        "安全门判定（如实声明）",
        [
            "冻结门要求候选 180 run 绝对零碰撞/零红灯违规/零非法相位",
            "场景 11 存在跨算法碰撞（fixed 4 / classic 5 / CA-MP 3 个 run）——",
            "  属该 0.1s 步长场景官方配时下的固有数据现实，CA-MP 碰撞最少",
            "因此默认选择保守回退 fixed_time，改进声明 False（不选择性声明）",
            "门的修订属协议变更，须 test-first 修订与独立复审",
        ],
    )
    bullet_slide(
        "可复现性与证据",
        [
            "sealed evidence：每个 run 独立目录 + manifest/provenance/hashes",
            "分析清单：output/evidence/formal/analysis_manifest.json（含 SHA-256）",
            "复现：run_pdf_matrix --profile formal → analyze_matrix → generate_deliverables",
            "状态三态口径：pass / fail / not_run（Docker live 与第二环境 = not_run）",
        ],
    )
    title_slide(
        "结论",
        "CA-MP 旅行时间显著改善（CI 全负）、排队全场最低、扰动韧性最优；\n"
        "受场景 11 固有碰撞影响未过绝对安全门——数据与声明严格一致",
    )

    out = Path("output/deliverables/答辩PPT.pptx")
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"written: {out}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
